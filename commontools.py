#!/usr/bin/env python
# _*_ coding: utf-8 _*_
# @Time : 2021/12/22 0022 2:26
# @Author : Every spring Xu ( Mr.Xu Fengchun)
# @Email :<thomsenw005@gmail.com>
# @File : dd.py
# @desc :

#自定义函数  自己写的函数库
#公共函数库
import os

from pyefun import *


#替换字符串特殊字符 并替换成Yingge Ceramica Park exo:YinggeCeramicaPark 形式
def word_str_replace(str):
    str=str.replace('(', '').replace(')', '').replace(' ', '').replace('&','')
    return str

#替换字符串特殊字符 并替换成 exo:Yingge-Ceramica-Park 形式
#https://www.itranslater.com/qa/details/2125465650226267136
def word_str_replace_two(str):
    '''
    :param str:
    :return:
    '''
    str=str.lower()  #把字母转换为小写
    str = str.replace('(', '').replace(')', '').replace(' & ', '').replace('&', '').replace(' / ','').replace(', ','').replace(' + ','')
    str=str.replace(' #','').replace(' - ','').replace('.','').replace("'",'').replace('&','').replace(' + ','').replace('/','').replace(': ','')
    str=str.replace('- ','').replace('_','').replace(' [','').replace('] ','').replace(' |','').replace('| ','')
    str=str.replace('[','').replace(']','').replace('! ','')
    new_str="-".join(str.split())

    return new_str

#md5加密
def md5Encode(str):
    import hashlib
    # 参数必须是byte类型，否则报Unicode-objects must be encoded before hashing错误
    m = hashlib.md5(str.encode(encoding='utf-8'))
    return m.hexdigest()


#https://stackoverflow.com/questions/62218673/how-to-extract-the-text-in-the-textarea-frame-of-the-deepl-page
# text="free download psd"
# c=deepl_request_post(text,'auto','ES')
# print(c)
def deepl_request_post(text,source_language,target_language):
    import requests
    import time
    url = "https://www2.deepl.com/jsonrpc"
    #text = "Hello, how are you today?"
    r = requests.post(
        url,
        json = {
            "jsonrpc":"2.0",
            "method": "LMT_handle_jobs",
            "params": {
                "jobs":[{
                    "kind":"default",
                    "raw_en_sentence": text,
                    "raw_en_context_before":[],
                    "raw_en_context_after":[],
                    "preferred_num_beams":4,
                    "quality":"fast"
                }],
                "lang":{
                    "user_preferred_langs":[source_language,target_language],
                    #"source_lang_user_selected":"auto",
                    "source_lang_user_selected": source_language,
                    "target_lang":target_language
                },
                "priority":-1,
                "commonJobParams":{},
                "timestamp": int(round(time.time() * 1000))
            },
            "id": 40890008
        }
    )
    result=r.json()
    #print(r.json())
    translations=result['result']['translations']
    #print(translations)
    for a in translations:
        #print(a['beams'])
        beams_list=a['beams']
    # for b in a['beams']:
    #     print(b['postprocessed_sentence'])

    #print(a['beams'][0]['postprocessed_sentence'])
    resutlt=beams_list[0]['postprocessed_sentence']
    return resutlt



# print(word_str_replace_two('LETTRRPRESS Photoshop LATYer Style & PSD '))
# print(md5Encode('a bad file name (really)'))




# 替换字符串
def spaceReplace(str):

    str=str.replace(' - ','')
    str=str.replace(': ','')
    str=str.replace('?','')
    str=str.replace(' & ','')
    str=str.replace('/','')
    str=str.replace(' ','-')
    str=str.replace("'",'')

    return str


# 2021/1/11 add
# https://blog.csdn.net/s1164548515/article/details/91430506
#获取文件夹file_dir下指定类型file_type的所有文件名
def file_name(file_dir,file_type=''):
    lst=[]
    for root,dirs,files in os.walk(file_dir):
        for file in files:
            if(file_type==''):
                lst.append(file)
            else:
                if os.path.splitext(file)[1]==str(file_type): #获取指定类型的文件名
                    lst.append(file)
    return  lst


#删除PPT最后一页 [有这样一个需求，就是现在我有大量的PPT文件，需要去掉每个PPT文件中的最后一页。]
# pip install -i https://pypi.tuna.tsinghua.edu.cn/simple python-pptx
from pptx import Presentation
def remove_ppt_last_page(pptFile):
    # 读取ppt
    prs = Presentation(pptFile)
    # 删除最后一页
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]
    # 保存新的ppt
    prs.save(pptFile)
    print(f"{pptFile},文件处理成功")